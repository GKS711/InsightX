"""
InsightX 簡報 generator

執行：
    cd ppt && /path/to/.venv/bin/python generate_v6.py

產出：
    ppt/InsightX.pptx

設計：
- 9 張 slide，magazine 風 (cream paper + ink + coral + forest)
- 16:9 widescreen
- 不是版本進化敘事，是「InsightX 是什麼 / 解什麼問題」的 portfolio 介紹
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pathlib import Path

# ─── Palette ────────────────────────────────────────────
INK       = RGBColor(0x1a, 0x1f, 0x1c)
INK_2     = RGBColor(0x5b, 0x66, 0x61)
INK_3     = RGBColor(0x8a, 0x93, 0x8e)
INK_4     = RGBColor(0xb0, 0xb6, 0xb3)
RULE      = RGBColor(0xcd, 0xd1, 0xce)
PAPER     = RGBColor(0xfa, 0xf7, 0xf2)
PAPER_2   = RGBColor(0xf3, 0xee, 0xe5)
CORAL     = RGBColor(0xd6, 0x5a, 0x3a)
FOREST    = RGBColor(0x2c, 0x5f, 0x2d)
WHITE     = RGBColor(0xff, 0xff, 0xff)
BLACK     = RGBColor(0x0d, 0x0f, 0x0e)

SERIF = "Noto Serif TC"
SANS  = "Noto Sans TC"
MONO  = "Menlo"

WIDTH = 13.333
HEIGHT = 7.5

HERE = Path(__file__).resolve().parent
HERO_IMG = HERE.parent / "src" / "static" / "v2" / "assets" / "hero-listening.png"
GAME_IMG = HERE.parent / "docs" / "screenshots" / "game-question.png"
LANDING_IMG = HERE.parent / "docs" / "screenshots" / "v4" / "01-landing.png"


# ─── Helpers ───────────────────────────────────────────
def add_bg(slide, color):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(WIDTH), Inches(HEIGHT))
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()


def add_text(slide, text, x, y, w, h, *, size=14, bold=False, color=INK, font=SANS,
             align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = valign
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = font
    return tb


def add_kicker(slide, num, label, x=0.5, y=0.4):
    add_text(slide, f"§ {num}", x, y, 1, 0.3, size=11, color=CORAL, font=MONO)
    add_text(slide, label.upper(), x + 0.7, y, 6, 0.3, size=11, color=INK_3, font=MONO)


def add_title(slide, text, x=0.5, y=0.85, w=12, h=1.2, color=INK):
    add_text(slide, text, x, y, w, h, size=40, bold=True, color=color, font=SERIF)


def add_subtitle(slide, text, x=0.5, y=2.0, w=12, h=0.6, color=INK_2):
    add_text(slide, text, x, y, w, h, size=15, color=color, font=SANS)


def add_rule_line(slide, x, y, w, color=RULE):
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Emu(9525))
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()


def add_image(slide, path, x, y, w, h):
    if not Path(path).exists():
        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        rect.fill.solid(); rect.fill.fore_color.rgb = PAPER_2
        rect.line.color.rgb = INK_4
        add_text(slide, f"(image missing: {Path(path).name})",
                 x, y + h/2 - 0.2, w, 0.4, size=11, color=INK_3, font=MONO, align=PP_ALIGN.CENTER)
        return
    slide.shapes.add_picture(str(path), Inches(x), Inches(y), Inches(w), Inches(h))


def card(slide, x, y, w, h, *, fill=PAPER_2, border=None):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    rect.fill.solid(); rect.fill.fore_color.rgb = fill
    if border is None:
        rect.line.fill.background()
    else:
        rect.line.color.rgb = border
    return rect


# ─── Build deck ─────────────────────────────────────────
prs = Presentation()
prs.slide_width = Inches(WIDTH)
prs.slide_height = Inches(HEIGHT)
blank = prs.slide_layouts[6]


# ╔══════════════════════════════════════════════════════════════╗
# ║ Slide 1: Cover                                                ║
# ╚══════════════════════════════════════════════════════════════╝
s = prs.slides.add_slide(blank)
add_bg(s, BLACK)
# Brand dot
dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.5), Inches(0.5), Inches(0.45), Inches(0.45))
dot.fill.solid(); dot.fill.fore_color.rgb = CORAL; dot.line.fill.background()
add_text(s, "i", 0.5, 0.43, 0.45, 0.45, size=22, bold=True, color=WHITE,
         font=SERIF, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
add_text(s, "insightx", 1.05, 0.5, 3, 0.5, size=20, bold=True, color=WHITE, font=SERIF)
add_text(s, "A I   顧 客 意 見 洞 察 平 台   ·   L I V E", 8.0, 0.55, 5, 0.3,
         size=10, color=INK_3, font=MONO, align=PP_ALIGN.RIGHT)

# Top kicker
add_text(s, "INSIGHTX · PROJECT INTRO · 2026", 0.7, 2.5, 8, 0.3,
         size=11, color=CORAL, font=MONO)

# Hero text — 兩行
add_text(s, "讀完所有評論，", 0.7, 3.0, 12, 1.4, size=72, bold=True, color=WHITE, font=SERIF)
add_text(s, "再給你一份報告。", 0.7, 4.4, 12, 1.4, size=72, bold=True, color=WHITE, font=SERIF)

# Subtitle
add_text(s,
         "輸入一個 Google Maps 或 YouTube 連結 ── AI 自動整理成情緒分析、SWOT、回覆草稿、行銷貼文、週行動計畫。",
         0.7, 6.0, 12, 0.6, size=15, color=INK_4, font=SANS)

# Footer rule + meta
add_rule_line(s, 0.7, 6.8, 11.9, color=RGBColor(0x44, 0x4d, 0x49))
add_text(s, "INSIGHTX  ·  MIT", 0.7, 6.95, 6, 0.3,
         size=10, color=INK_4, font=MONO)
add_text(s, "Jordan711-insightx-demo.hf.space", 8, 6.95, 5, 0.3,
         size=10, color=INK_4, font=MONO, align=PP_ALIGN.RIGHT)


# ╔══════════════════════════════════════════════════════════════╗
# ║ Slide 2: 痛點                                                 ║
# ╚══════════════════════════════════════════════════════════════╝
s = prs.slides.add_slide(blank)
add_bg(s, PAPER)
add_kicker(s, "01", "痛點")
add_title(s, "顧客每天在留言，但意見很少變成決策。")
add_subtitle(s, "管理者每天打開後台看到一堆五星評論，但下一步要做什麼？沒人說得出來。")

pains = [
    ("太碎了", "意見散在 Google Maps、YouTube、IG、LINE 各處，\n沒人有時間全部讀完。"),
    ("太多了", "50 則留言已經是負擔，500 則直接放棄。看得到「很多人說\n好吃」，看不到「冷氣太強客人不想久坐」這種真正的細節。"),
    ("回應壓力大", "負評上線，店長被情緒帶著走，回覆要嘛太硬要嘛太軟，\n反而失去挽回客人的機會。"),
    ("新店長沒練習場", "剛升上店長的人沒處理過真正的危機，\n學費全部是用真實客人賠的。"),
]
# 2x2 grid
positions = [(0.7, 3.0), (7.0, 3.0), (0.7, 5.1), (7.0, 5.1)]
for (title, body), (x, y) in zip(pains, positions):
    card(s, x, y, 5.7, 1.9, fill=PAPER_2)
    # coral accent bar
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.08), Inches(1.9))
    bar.fill.solid(); bar.fill.fore_color.rgb = CORAL; bar.line.fill.background()
    add_text(s, title, x + 0.3, y + 0.25, 5.2, 0.5, size=20, bold=True, color=INK, font=SERIF)
    add_text(s, body, x + 0.3, y + 0.85, 5.2, 1.0, size=13, color=INK_2, font=SANS)


# ╔══════════════════════════════════════════════════════════════╗
# ║ Slide 3: 解決方案                                              ║
# ╚══════════════════════════════════════════════════════════════╝
s = prs.slides.add_slide(blank)
add_bg(s, PAPER)
add_kicker(s, "02", "解決方案")
add_title(s, "一個 URL 換一份完整報告。")
add_subtitle(s, "把「讀評論 → 想策略 → 寫回覆 → 練手」串成同一條流水線。")

bullets = [
    ("跨平台抓料，不用開瀏覽器",
     "Google Maps 用 Serper API 抓店家評論，YouTube 用官方 Data API 抓影片留言。\n速度快、也不容易被反爬擋掉。"),
    ("9 個 AI 功能共用一份資料",
     "情緒分析、SWOT、回覆草稿、行銷文案、根源分析、週計畫、培訓劇本、內部信、AI 對話顧問。\n依平台切換語氣（餐廳 / 零售 / YouTuber）。"),
    ("多店家工作區",
     "每位使用者有自己的工作區，可以新增多個店家、保留完整歷史紀錄。\n隨時回看上個月對某家分店的分析。"),
    ("管理者決策模擬遊戲",
     "把真實負評倒進小遊戲，AI 當你的虛擬顧問即時給回饋。\n練手不用拿真客人試刀。"),
]
for i, (title, body) in enumerate(bullets):
    y = 2.85 + i * 1.15
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(y + 0.05), Inches(0.05), Inches(1.0))
    bar.fill.solid(); bar.fill.fore_color.rgb = FOREST; bar.line.fill.background()
    add_text(s, title, 0.95, y, 11, 0.45, size=18, bold=True, color=INK, font=SERIF)
    add_text(s, body, 0.95, y + 0.5, 11.5, 0.7, size=13, color=INK_2, font=SANS)


# ╔══════════════════════════════════════════════════════════════╗
# ║ Slide 4: 9 個 AI 功能                                          ║
# ╚══════════════════════════════════════════════════════════════╝
s = prs.slides.add_slide(blank)
add_bg(s, PAPER)
add_kicker(s, "03", "ai 功能")
add_title(s, "9 個 AI 功能，共用一份原始評論。")
add_subtitle(s, "依平台切換語氣 ── 餐廳老闆看到的是門市建議，YouTuber 看到的是頻道策略。")

functions = [
    ("01", "情緒分析", "正負向 + 主題分布"),
    ("02", "SWOT", "自動生成戰略矩陣"),
    ("03", "回覆草稿", "每則負評一個對應草稿"),
    ("04", "行銷文案", "門市活動 / 影片宣傳"),
    ("05", "根源分析", "找出真正的痛點"),
    ("06", "週行動計畫", "一週要做哪些具體事項"),
    ("07", "培訓劇本", "員工 / 剪輯師訓練教材"),
    ("08", "內部信", "門市 / 團隊週報"),
    ("09", "AI 顧問", "隨時可問的虛擬導師"),
]
# 3x3 grid
for i, (num, title, body) in enumerate(functions):
    row = i // 3
    col = i % 3
    x = 0.7 + col * 4.2
    y = 3.0 + row * 1.4
    card(s, x, y, 3.9, 1.2, fill=PAPER_2)
    add_text(s, num, x + 0.25, y + 0.2, 1, 0.3, size=10, color=CORAL, font=MONO)
    add_text(s, title, x + 0.25, y + 0.45, 3.5, 0.4, size=18, bold=True, color=INK, font=SERIF)
    add_text(s, body, x + 0.25, y + 0.85, 3.5, 0.4, size=12, color=INK_2, font=SANS)


# ╔══════════════════════════════════════════════════════════════╗
# ║ Slide 5: 管理者決策模擬                                        ║
# ╚══════════════════════════════════════════════════════════════╝
s = prs.slides.add_slide(blank)
add_bg(s, PAPER)
add_kicker(s, "04", "決策模擬")
add_title(s, "用真實負評，練決策。")
add_subtitle(s, "新任店長最痛的事 ── 第一次處理客訴。這個遊戲讓你在真實負評上練手。")

# Left: game screenshot
add_image(s, GAME_IMG, 0.7, 3.1, 6.0, 3.8)

# Right: 3-step explanation
steps = [
    ("1.", "AI 出題", "把真實負評變成情境問句"),
    ("2.", "你選回應", "從多個策略選一個你會做的"),
    ("3.", "AI 給回饋", "評估你的選擇，給情商分數 + 建議"),
]
for i, (num, title, body) in enumerate(steps):
    y = 3.2 + i * 1.1
    add_text(s, num, 7.2, y, 0.6, 0.5, size=28, bold=True, color=CORAL, font=SERIF)
    add_text(s, title, 7.9, y, 5, 0.45, size=20, bold=True, color=INK, font=SERIF)
    add_text(s, body, 7.9, y + 0.5, 5, 0.45, size=13, color=INK_2, font=SANS)

# Bottom note
add_text(s, "── 學費，不用拿真客人賠。", 7.2, 6.6, 6, 0.4, size=14, color=FOREST, font=SERIF)


# ╔══════════════════════════════════════════════════════════════╗
# ║ Slide 6: 3 個技術選擇 (白話)                                   ║
# ╚══════════════════════════════════════════════════════════════╝
s = prs.slides.add_slide(blank)
add_bg(s, PAPER)
add_kicker(s, "05", "技術選擇")
add_title(s, "3 個決定，影響使用體驗。")
add_subtitle(s, "不堆技術名詞 ── 講為什麼這樣做、解了什麼問題。")

tech_choices = [
    ("多人也能安全使用",
     "最初版本所有人共用一個帳號，任何人打開都看到別人的資料。\n"
     "現在每位訪客一進站，自動拿到一個身分（用 cookie 記下來），\n"
     "所有資料綁在這身分上。沒有註冊、沒有密碼，但每個人的工作區完全分開。"),
    ("不會因為 AI 抖一下就掛",
     "免費版 Gemini 半夜常常會卡。我做了一個自動換模型的機制：\n"
     "主用快的，失敗就降到大的，再不行換 Google 的旗艦，最後 lite 版兜底。\n"
     "同一次請求一路 fallback，使用者完全感覺不到。"),
    ("雙 AI 寫 code",
     "改動比較大的時候，我習慣讓 Codex（另一個 AI 助手）幫我看一遍程式碼。\n"
     "一個 AI 寫、另一個 AI 挑毛病。抓出來的問題比自己看更多 ──\n"
     "等於多了一個免費的 reviewer。"),
]
for i, (title, body) in enumerate(tech_choices):
    y = 3.0 + i * 1.45
    # Small green dot
    d = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.7), Inches(y + 0.18), Inches(0.2), Inches(0.2))
    d.fill.solid(); d.fill.fore_color.rgb = FOREST; d.line.fill.background()
    add_text(s, title, 1.1, y, 11, 0.5, size=20, bold=True, color=INK, font=SERIF)
    add_text(s, body, 1.1, y + 0.55, 11.5, 0.9, size=13, color=INK_2, font=SANS)


# ╔══════════════════════════════════════════════════════════════╗
# ║ Slide 7: 技術棧 + 系統架構                                     ║
# ╚══════════════════════════════════════════════════════════════╝
s = prs.slides.add_slide(blank)
add_bg(s, PAPER)
add_kicker(s, "06", "系統架構")
add_title(s, "從一張圖看整個系統。")
add_subtitle(s, "後端 FastAPI，AI 是 Gemini，部署在 Hugging Face Spaces ── 全部用免費方案組起來。")

# Left: tech stack
add_text(s, "TECH STACK", 0.7, 3.0, 5, 0.3, size=11, color=CORAL, font=MONO)
stack_items = [
    ("後端", "FastAPI · Python 3.10+"),
    ("資料庫", "Turso · SQLite (libsql)"),
    ("AI", "Google Gemini (多模型 fallback)"),
    ("爬蟲", "Serper API + YouTube Data API"),
    ("前端", "React 18 + Tailwind CSS"),
    ("部署", "Docker on HF Spaces ($0/月)"),
]
for i, (k, v) in enumerate(stack_items):
    y = 3.45 + i * 0.55
    add_text(s, k, 0.7, y, 1.5, 0.4, size=13, bold=True, color=INK, font=SANS)
    add_text(s, v, 2.2, y, 4.5, 0.4, size=13, color=INK_2, font=MONO)

# Right: architecture diagram (boxes + arrows)
arch_x = 7.3
arch_y = 3.0
add_text(s, "DATA FLOW", arch_x, arch_y, 5, 0.3, size=11, color=CORAL, font=MONO)

def box(x, y, w, h, label, *, fill=PAPER_2, color=INK, font_size=12, bold=True):
    rect = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    rect.fill.solid(); rect.fill.fore_color.rgb = fill
    rect.line.color.rgb = INK_4
    add_text(s, label, x, y, w, h, size=font_size, bold=bold, color=color, font=SANS,
             align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

def arrow(x, y, w, h):
    add_text(s, "↓", x, y, w, h, size=14, color=INK_3, font=SANS, align=PP_ALIGN.CENTER)

# Browser
box(arch_x, arch_y + 0.5, 5.3, 0.5, "使用者瀏覽器", fill=PAPER)
arrow(arch_x, arch_y + 1.05, 5.3, 0.3)
# FastAPI
box(arch_x, arch_y + 1.4, 5.3, 0.5, "FastAPI 後端", fill=PAPER_2)
# v4 / v5 split
box(arch_x, arch_y + 2.0, 2.5, 0.45, "/api/v4 (一次性)", fill=PAPER, font_size=11)
box(arch_x + 2.8, arch_y + 2.0, 2.5, 0.45, "/api/v5 (持久工作區)", fill=PAPER, font_size=11)
arrow(arch_x, arch_y + 2.5, 5.3, 0.3)
# 4 externals
ext_y = arch_y + 2.95
ext_w = 1.25
ext_gap = 0.05
for i, (label, color) in enumerate([
    ("Gemini AI", FOREST),
    ("Turso DB", CORAL),
    ("Serper API", INK_2),
    ("YouTube API", INK_2),
]):
    x = arch_x + i * (ext_w + ext_gap)
    box(x, ext_y, ext_w, 0.55, label, fill=PAPER_2, color=color, font_size=11)


# ╔══════════════════════════════════════════════════════════════╗
# ║ Slide 8: 開發過程學到的                                        ║
# ╚══════════════════════════════════════════════════════════════╝
s = prs.slides.add_slide(blank)
add_bg(s, PAPER)
add_kicker(s, "07", "學習心得")
add_title(s, "做這個專案學到的事。")
add_subtitle(s, "從第一版 demo 到能上線給人用的版本，學到幾件之前沒想過的事。")

lessons = [
    ("早期選錯路徑代價很大",
     "最初用 Playwright 自動瀏覽器抓 Google Maps，後來 Google 改了反爬機制只好整段砍掉換成 Serper API。\n"
     "如果一開始先試 API，會省下大概兩個月。"),
    ("API 抖動比 prompt 寫不好還麻煩",
     "免費版 Gemini 半夜會 500，再怎麼調 prompt 都救不了。做了多模型 fallback 之後，\n"
     "可用率從大概 85% 拉到 99%。"),
    ("同步反而比較好 debug",
     "原本後端是 async 寫的，後來為了搭 Turso 改回同步。原本以為是退步，\n"
     "結果發現用 threading 反而出問題比較好追，效能也沒掉。"),
    ("多一個 AI 幫忙看 code 真的有差",
     "自己寫自己看很容易漏掉邊角，找 Codex 一起 review 後抓到不少會炸的 bug ──\n"
     "多花一點時間，少踩很多坑。"),
]
for i, (title, body) in enumerate(lessons):
    y = 2.9 + i * 1.05
    add_text(s, f"0{i+1}", 0.7, y + 0.1, 0.8, 0.5, size=24, bold=True, color=CORAL, font=SERIF)
    add_text(s, title, 1.6, y, 11, 0.5, size=17, bold=True, color=INK, font=SERIF)
    add_text(s, body, 1.6, y + 0.5, 11.5, 0.7, size=12, color=INK_2, font=SANS)


# ╔══════════════════════════════════════════════════════════════╗
# ║ Slide 9: 線上試用 (closing, 黑底)                              ║
# ╚══════════════════════════════════════════════════════════════╝
s = prs.slides.add_slide(blank)
add_bg(s, BLACK)
dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.5), Inches(0.5), Inches(0.45), Inches(0.45))
dot.fill.solid(); dot.fill.fore_color.rgb = CORAL; dot.line.fill.background()
add_text(s, "i", 0.5, 0.43, 0.45, 0.45, size=22, bold=True, color=WHITE,
         font=SERIF, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
add_text(s, "insightx", 1.05, 0.5, 3, 0.5, size=20, bold=True, color=WHITE, font=SERIF)
add_text(s, "T H A N K   Y O U", 9.5, 0.55, 3.5, 0.3,
         size=10, color=INK_3, font=MONO, align=PP_ALIGN.RIGHT)

add_text(s, "現在就試。", 0.7, 2.5, 8, 0.5, size=11, color=CORAL, font=MONO)
add_text(s, "免註冊。", 0.7, 3.0, 12, 1.4, size=84, bold=True, color=WHITE, font=SERIF)
add_text(s, "免安裝。月費 $0。", 0.7, 4.4, 12, 1.4, size=72, bold=True, color=WHITE, font=SERIF)

add_text(s,
         "Demo 跑在 Hugging Face Spaces 免費方案。48 小時沒人用會睡眠，第一次點要等 30 秒醒過來。",
         0.7, 5.9, 12, 0.5, size=13, color=INK_4, font=SANS)

add_rule_line(s, 0.7, 6.5, 11.9, color=RGBColor(0x44, 0x4d, 0x49))
add_text(s, "🌐  線上 Demo", 0.7, 6.7, 4, 0.3, size=11, color=INK_3, font=MONO)
add_text(s, "Jordan711-insightx-demo.hf.space", 0.7, 7.0, 6, 0.4,
         size=14, bold=True, color=WHITE, font=MONO)
add_text(s, "📦  GitHub", 7.5, 6.7, 4, 0.3, size=11, color=INK_3, font=MONO)
add_text(s, "github.com/GKS711/InsightX", 7.5, 7.0, 6, 0.4,
         size=14, bold=True, color=WHITE, font=MONO)


# ─── Save ───────────────────────────────────────────────
output_path = HERE / "InsightX.pptx"
prs.save(str(output_path))
print(f"✓ Saved {output_path} ({len(prs.slides)} slides)")
