"""
InsightX v4.0.0 PPT Generator

執行：
    cd ppt && python3 generate.py

產出：
    ppt/InsightX_v4.0.0.pptx

設計：
- 配色：與 v4 UI 一致
    背景 cream  #faf7f2
    ink     #1a1f1c (主文字)
    ink-2   #5b6661 (次文字)
    ink-3   #8a938e (caption)
    rule    #cdd1ce (淡分隔線)
    coral   #d65a3a (accent)
    forest  #2c5f2d (positive)
- 字體：標題 serif / 內文 sans
- 結構：14 張，sandwich 結構（dark 封面/結尾，light 內容）
- 視覺一致：每張左上角 § 編號 + kicker（仿 v4 dashboard 章節風）
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pathlib import Path

# ─── Palette ────────────────────────────────────────────
INK       = RGBColor(0x1a, 0x1f, 0x1c)
INK_2     = RGBColor(0x5b, 0x66, 0x61)
INK_3     = RGBColor(0x8a, 0x93, 0x8e)
INK_4     = RGBColor(0xb0, 0xb6, 0xb3)
RULE      = RGBColor(0xcd, 0xd1, 0xce)
PAPER     = RGBColor(0xfa, 0xf7, 0xf2)
PAPER_2   = RGBColor(0xf3, 0xee, 0xe5)
CORAL     = RGBColor(0xd6, 0x5a, 0x3a)
FOREST    = RGBColor(0x2c, 0x5f, 0x2d)
WHITE     = RGBColor(0xff, 0xff, 0xff)
BLACK     = RGBColor(0x0d, 0x0f, 0x0e)

# Fonts
SERIF = "Noto Serif TC"   # 中英都好看
SANS  = "Noto Sans TC"
MONO  = "Menlo"

# Layout: 16:9 @ 13.333 x 7.5 in (LAYOUT_WIDE)
WIDTH = 13.333
HEIGHT = 7.5

# ─── Resolve assets path ───────────────────────────────
HERE = Path(__file__).resolve().parent
ASSETS = HERE / "assets"
SCREENSHOTS_SRC = HERE.parent / "docs" / "screenshots" / "v4"

# ─── Helpers ───────────────────────────────────────────
def add_bg(slide, color):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(WIDTH), Inches(HEIGHT))
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    return bg

def add_text(slide, text, x, y, w, h, *, size=14, bold=False, color=INK, font=SANS,
             align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, italic=False):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.word_wrap = True
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    f = r.font
    f.name = font
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.color.rgb = color
    return tb

def add_kicker(slide, num, label, x=0.5, y=0.4):
    """左上 §0X · KICKER style mono caption"""
    add_text(slide, f"§{num}    {label.upper()}", x, y, 6, 0.3,
             size=10, font=MONO, color=CORAL, bold=False)

def add_title(slide, text, x=0.5, y=0.85, w=12, h=1.2, color=INK):
    add_text(slide, text, x, y, w, h, size=44, bold=True, color=color, font=SERIF)

def add_subtitle(slide, text, x=0.5, y=2.1, w=12, h=0.5, color=INK_2):
    add_text(slide, text, x, y, w, h, size=18, color=color, font=SANS)

def add_rule_line(slide, x, y, w, color=RULE):
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.012))
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()
    return line

def add_image(slide, path, x, y, w, h):
    if Path(path).exists():
        slide.shapes.add_picture(str(path), Inches(x), Inches(y), Inches(w), Inches(h))
    else:
        # placeholder
        ph = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        ph.fill.solid()
        ph.fill.fore_color.rgb = PAPER_2
        ph.line.color.rgb = RULE
        ph.line.width = Pt(0.5)
        add_text(slide, f"[ {Path(path).name} ]", x, y + h/2 - 0.2, w, 0.4, size=11,
                 color=INK_3, font=MONO, align=PP_ALIGN.CENTER)

# ─── Build deck ─────────────────────────────────────────
prs = Presentation()
prs.slide_width = Inches(WIDTH)
prs.slide_height = Inches(HEIGHT)
blank = prs.slide_layouts[6]  # blank layout

# ╔══════════════════════════════════════════════════════════════╗
# ║ Slide 1: Cover                                                ║
# ╚══════════════════════════════════════════════════════════════╝
s = prs.slides.add_slide(blank)
add_bg(s, BLACK)
# logo dot
dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.5), Inches(0.5), Inches(0.45), Inches(0.45))
dot.fill.solid(); dot.fill.fore_color.rgb = CORAL; dot.line.fill.background()
add_text(s, "i", 0.5, 0.43, 0.45, 0.45, size=22, bold=True, color=WHITE,
         font=SERIF, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
add_text(s, "insightx", 1.05, 0.5, 3, 0.5, size=20, bold=True, color=WHITE, font=SERIF)
# version pill (upper right)
add_text(s, "V 4 . 0 . 0   ·   L I V E", 10.5, 0.55, 2.5, 0.3,
         size=10, color=INK_3, font=MONO, align=PP_ALIGN.RIGHT)

# Hero typography
add_text(s, "AI · SHOP / CHANNEL INTELLIGENCE", 0.7, 2.5, 8, 0.3,
         size=11, color=CORAL, font=MONO)
add_text(s, "把顧客每一句話，", 0.7, 3.0, 12, 1.4, size=72, bold=True, color=WHITE, font=SERIF)
add_text(s, "讀給老闆聽。", 0.7, 4.4, 12, 1.4, size=72, bold=True, color=WHITE, font=SERIF)
add_text(s, "AI 驅動的雙平台顧客回饋分析 — Google Maps 店家評論 + YouTube 影片留言",
         0.7, 6.0, 12, 0.5, size=15, color=INK_4, font=SANS)
# bottom rule
add_rule_line(s, 0.7, 6.8, 11.9, color=RGBColor(0x44, 0x4d, 0x49))
add_text(s, "INSIGHTX  ·  AI ADVISOR REPORT  ·  MIT", 0.7, 6.95, 12, 0.3,
         size=10, color=INK_4, font=MONO)
add_text(s, "github.com/GKS711/InsightX", 8.5, 6.95, 4.5, 0.3,
         size=10, color=INK_4, font=MONO, align=PP_ALIGN.RIGHT)


# ╔══════════════════════════════════════════════════════════════╗
# ║ Slide 2: 痛點                                                  ║
# ╚══════════════════════════════════════════════════════════════╝
s = prs.slides.add_slide(blank)
add_bg(s, PAPER)
add_kicker(s, "01", "the problem")
add_title(s, "讀不完，是死循環。")
add_subtitle(s, "店家有 100 則評論、創作者有 1000 則留言 — 沒人有時間逐條讀。")

# 兩欄痛點
def pain_column(x, badge, badge_label, headline, lines):
    # badge
    b = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(3.0), Inches(1.3), Inches(0.3))
    b.fill.solid(); b.fill.fore_color.rgb = INK; b.line.fill.background()
    add_text(s, badge, x, 3.0, 1.3, 0.3, size=11, color=WHITE, font=MONO,
             align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_text(s, badge_label, x + 1.5, 3.02, 4, 0.3, size=11, color=INK_3, font=MONO)
    # headline
    add_text(s, headline, x, 3.5, 5.5, 0.7, size=24, bold=True, color=INK, font=SERIF)
    # lines
    for i, ln in enumerate(lines):
        add_text(s, ln, x, 4.5 + i * 0.55, 5.5, 0.5, size=14, color=INK_2, font=SANS)

pain_column(0.7, "G", "GOOGLE 店家", "讀完 200 則評論、整理回覆？",
            ["平均一則 30-60 秒，整月評論 200 則 = 整天 only this。",
             "沒空讀 → 沒洞察 → 看不出顧客真正在意什麼。",
             "回覆草稿要寫到讓客戶感受到誠意又不卑微，難。"])
pain_column(7.0, "Y", "YOUTUBE 創作者", "上千則留言、抓不到改進方向？",
            ["熱門影片留言破千，下一支該往哪個方向做？看不出。",
             "正評負評混在一起，演算法到底為什麼沒推？沒線索。",
             "想跟頻道資料對話，需要一個讀過所有留言的副駕駛。"])


# ╔══════════════════════════════════════════════════════════════╗
# ║ Slide 3: Solution overview                                    ║
# ╚══════════════════════════════════════════════════════════════╝
s = prs.slides.add_slide(blank)
add_bg(s, PAPER)
add_kicker(s, "02", "the solution")
add_title(s, "一個 URL，一份顧問報告。")
add_subtitle(s, "貼上連結 → AI 自動爬取 + 分析 + 生成 9 種策略產出 → 報告秒回。")

# 3-step horizontal flow
steps = [
    ("01", "貼上 URL", "Google Maps 店家網址\n或 YouTube 影片網址"),
    ("02", "AI 分析", "Serper / YouTube API 爬留言\nGemini gemma-4-31b 分析"),
    ("03", "完整報告", "Sentiment / SWOT / Reply\nMarketing / Plan / Chat"),
]
for i, (num, title, desc) in enumerate(steps):
    x = 0.7 + i * 4.2
    # number
    add_text(s, num, x, 3.4, 1.2, 1.2, size=72, bold=True, color=CORAL, font=SERIF)
    add_text(s, title, x + 1.5, 3.6, 2.5, 0.5, size=22, bold=True, color=INK, font=SERIF)
    add_text(s, desc, x + 1.5, 4.2, 2.5, 1.2, size=12, color=INK_2, font=SANS)
    # arrow
    if i < 2:
        add_text(s, "→", x + 3.7, 3.7, 0.5, 0.6, size=24, color=INK_3, font=SANS)

# bottom callout
add_rule_line(s, 0.7, 5.9, 11.9)
add_text(s, "30 SECONDS", 0.7, 6.2, 4, 0.3, size=11, color=INK_3, font=MONO)
add_text(s, "從 URL 到完整 dashboard，只要 30 秒（不含 LLM 生成時間）。",
         0.7, 6.5, 12, 0.5, size=14, color=INK_2, font=SANS)


# ╔══════════════════════════════════════════════════════════════╗
# ║ Slide 4: Two platforms                                        ║
# ╚══════════════════════════════════════════════════════════════╝
s = prs.slides.add_slide(blank)
add_bg(s, PAPER)
add_kicker(s, "03", "two platforms")
add_title(s, "兩種來源，同一份顧問。")
add_subtitle(s, "店家評論看顧客「怎麼想」；YouTube 留言看外人「怎麼看」。")

# Two cards
def platform_card(x, badge_color, brand_letter, title, sub, bullets):
    card = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(3.0), Inches(5.7), Inches(3.7))
    card.fill.solid(); card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = INK; card.line.width = Pt(1)
    # top accent strip
    strip = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(3.0), Inches(5.7), Inches(0.08))
    strip.fill.solid(); strip.fill.fore_color.rgb = badge_color; strip.line.fill.background()
    # brand square
    sq = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x + 0.4), Inches(3.4), Inches(0.7), Inches(0.7))
    sq.fill.solid(); sq.fill.fore_color.rgb = badge_color; sq.line.fill.background()
    add_text(s, brand_letter, x + 0.4, 3.4, 0.7, 0.7, size=28, bold=True, color=WHITE,
             font=SERIF, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    # mono kicker
    add_text(s, sub, x + 1.3, 3.45, 4, 0.3, size=10, color=INK_3, font=MONO)
    add_text(s, title, x + 1.3, 3.7, 4, 0.5, size=22, bold=True, color=INK, font=SERIF)
    # bullets
    for i, b in enumerate(bullets):
        add_text(s, "▸", x + 0.4, 4.7 + i * 0.55, 0.3, 0.5, size=14, color=CORAL, font=SANS)
        add_text(s, b, x + 0.7, 4.7 + i * 0.55, 5, 0.5, size=13, color=INK_2, font=SANS)

platform_card(0.7, RGBColor(0x42, 0x85, 0xF4), "G", "Google 評論",
              "REVIEWS · STRUCTURED",
              ["Serper /maps + /reviews 純 API",
               "拿到星等、文字、作者、時間",
               "可分頁抓 ratingCount > 100 的店家",
               "適合餐廳、零售、服務業"])
platform_card(6.9, RGBColor(0xFF, 0x00, 0x00), "Y", "YouTube 影片",
              "COMMENTS · UNSTRUCTURED",
              ["YouTube Data API v3（主）",
               "youtube-comment-downloader（備用）",
               "拿到留言文字、讚數、作者",
               "適合創作者、頻道成長診斷"])


# ╔══════════════════════════════════════════════════════════════╗
# ║ Slide 5: Dashboard hero (screenshot)                          ║
# ╚══════════════════════════════════════════════════════════════╝
s = prs.slides.add_slide(blank)
add_bg(s, PAPER)
add_kicker(s, "04", "dashboard · §01 hero")
add_title(s, "一眼看完核心指標。")
add_subtitle(s, "店名、目前評分、90 天趨勢、情感分布 — 不滑動就看到。")
add_image(s, ASSETS / "04-hero.png", 0.7, 3.0, 11.9, 4.0)


# ╔══════════════════════════════════════════════════════════════╗
# ║ Slide 6: §02 themes                                           ║
# ╚══════════════════════════════════════════════════════════════╝
s = prs.slides.add_slide(blank)
add_bg(s, PAPER)
add_kicker(s, "05", "dashboard · §02 themes")
add_title(s, "他們實際在談的，不是你以為的。")
add_subtitle(s, "正向 / 負面 Top 3 主題 + 真實百分比 — 沒抽到主題就空著，不灌罐頭數字。")
add_image(s, ASSETS / "05-themes.png", 0.7, 3.0, 11.9, 4.0)


# ╔══════════════════════════════════════════════════════════════╗
# ║ Slide 7: §03 SWOT                                             ║
# ╚══════════════════════════════════════════════════════════════╝
s = prs.slides.add_slide(blank)
add_bg(s, PAPER)
add_kicker(s, "06", "dashboard · §03 swot")
add_title(s, "你站在哪裡，你能去哪裡。")
add_subtitle(s, "SWOT 四象限，每條都標 evidence-backed + 引用觸發它的評論百分比。")
add_image(s, ASSETS / "06-swot.png", 0.7, 3.0, 11.9, 4.0)


# ╔══════════════════════════════════════════════════════════════╗
# ║ Slide 8: §04 reviews                                          ║
# ╚══════════════════════════════════════════════════════════════╝
s = prs.slides.add_slide(blank)
add_bg(s, PAPER)
add_kicker(s, "07", "dashboard · §04 original material")
add_title(s, "永遠看得見來源。")
add_subtitle(s, "最多 50 則原始評論／留言 + 誠實 caption「分析了 N 則 · 顯示精選的 50 則」。")
add_image(s, ASSETS / "07-reviews.png", 0.7, 3.0, 11.9, 4.0)


# ╔══════════════════════════════════════════════════════════════╗
# ║ Slide 9: §07 toolbox - 5 generators                           ║
# ╚══════════════════════════════════════════════════════════════╝
s = prs.slides.add_slide(blank)
add_bg(s, PAPER)
add_kicker(s, "08", "toolbox · 5 generators")
add_title(s, "從洞察到行動，五件可以今天就做的事。")
add_subtitle(s, "工具箱：負評回覆草稿、行銷文案、週行動計畫、員工培訓、內部公告信。")

# 兩欄：左截圖、右說明
add_image(s, ASSETS / "08-week-plan.png", 0.7, 3.0, 6.5, 4.0)
tools = [
    ("負評回覆草稿", "對個別痛點回，附自我審查面板"),
    ("行銷文案", "對齊真實優勢，社群可貼"),
    ("週行動計畫", "7 天 to-do（負責人+預期結果）"),
    ("員工培訓劇本", "針對最大負面主題做話術"),
    ("內部公告信", "管理層／員工同步認知"),
]
for i, (name, desc) in enumerate(tools):
    y = 3.1 + i * 0.78
    # number circle
    n = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.5), Inches(y), Inches(0.4), Inches(0.4))
    n.fill.solid(); n.fill.fore_color.rgb = INK; n.line.fill.background()
    add_text(s, str(i+1), 7.5, y, 0.4, 0.4, size=12, bold=True, color=WHITE,
             font=SANS, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_text(s, name, 8.05, y - 0.05, 4.5, 0.4, size=15, bold=True, color=INK, font=SERIF)
    add_text(s, desc, 8.05, y + 0.32, 4.5, 0.4, size=11, color=INK_2, font=SANS)


# ╔══════════════════════════════════════════════════════════════╗
# ║ Slide 10: §AI Advisor                                         ║
# ╚══════════════════════════════════════════════════════════════╝
s = prs.slides.add_slide(blank)
add_bg(s, PAPER)
add_kicker(s, "09", "ai advisor · chat")
add_title(s, "問它關於你店的事。")
add_subtitle(s, "AI 顧問的 context 只有你的資料 — 不是通用 ChatGPT。")
add_image(s, ASSETS / "10-ai-advisor.png", 0.7, 3.0, 11.9, 4.0)


# ╔══════════════════════════════════════════════════════════════╗
# ║ Slide 11: Tech architecture                                   ║
# ╚══════════════════════════════════════════════════════════════╝
s = prs.slides.add_slide(blank)
add_bg(s, PAPER)
add_kicker(s, "10", "architecture")
add_title(s, "零瀏覽器，全 API。")
add_subtitle(s, "FastAPI + React 18 single-file SPA + Gemini + Serper / YouTube Data API。")

stack = [
    ("FRONTEND", "React 18 UMD + @babel/standalone\nsingle-file SPA · 不需 build step"),
    ("BACKEND",  "FastAPI · Python 3.10+\n結構化 SSE /api/v4/analyze-stream"),
    ("AI",       "Google Gemini gemma-4-31b-it\n9 個 platform-aware feature endpoint"),
    ("SCRAPERS", "Serper API（店家）\nYouTube Data API v3 + library fallback"),
    ("DEPLOY",   "Docker + Docker Compose\nMIT License · 開源"),
]
for i, (label, body) in enumerate(stack):
    y = 3.0 + i * 0.75
    add_text(s, label, 0.7, y, 1.6, 0.4, size=10, color=CORAL, font=MONO)
    add_text(s, body, 2.5, y - 0.05, 10, 0.7, size=12, color=INK, font=SANS)
    add_rule_line(s, 0.7, y + 0.6, 11.9)


# ╔══════════════════════════════════════════════════════════════╗
# ║ Slide 12: 4 invariants                                        ║
# ╚══════════════════════════════════════════════════════════════╝
s = prs.slides.add_slide(blank)
add_bg(s, PAPER)
add_kicker(s, "11", "engineering invariants")
add_title(s, "四條鎖定的工程紀律。")
add_subtitle(s, "跨後端 + 前端三層強制，任何後續修改都不能 regress。")

invariants = [
    ("①", "Frontend timeoutMs ≥ Backend total_timeout_s + 5s buffer"),
    ("②", "Service 層失敗一律 raise — 不回 fallback dict 造成 silent degradation"),
    ("③", "Retry 走 exception type — 不靠字串比對"),
    ("④", "Prompt 骨架對齊 <pre> renderer — 用 【】 ◆　▸ 純文字結構，禁 markdown"),
]
for i, (num, body) in enumerate(invariants):
    y = 3.2 + i * 0.85
    add_text(s, num, 0.7, y, 1.0, 0.7, size=44, bold=True, color=CORAL, font=SERIF)
    add_text(s, body, 1.7, y + 0.1, 11, 0.55, size=15, color=INK, font=SANS)


# ╔══════════════════════════════════════════════════════════════╗
# ║ Slide 13: Quick start                                         ║
# ╚══════════════════════════════════════════════════════════════╝
s = prs.slides.add_slide(blank)
add_bg(s, PAPER)
add_kicker(s, "12", "quick start")
add_title(s, "3 步驟，5 分鐘啟動。")

steps = [
    ("01", "Clone & 設定", "git clone https://github.com/GKS711/InsightX.git\ncp .env.example .env\n編輯 .env：GEMINI_API_KEY / SERPER_API_KEY / YOUTUBE_API_KEY"),
    ("02", "安裝 + 啟動", "pip install -r requirements.txt\npython -m uvicorn src.main:app --host 0.0.0.0 --port 8000"),
    ("03", "開瀏覽器", "http://localhost:8000\n貼上 Google Maps 或 YouTube 網址 → 開始分析"),
]
for i, (num, title, code) in enumerate(steps):
    y = 3.0 + i * 1.4
    add_text(s, num, 0.7, y, 1.0, 1.2, size=56, bold=True, color=CORAL, font=SERIF)
    add_text(s, title, 1.9, y + 0.15, 3.5, 0.5, size=20, bold=True, color=INK, font=SERIF)
    add_text(s, code, 5.5, y + 0.05, 7.5, 1.3, size=11, color=INK_2, font=MONO)


# ╔══════════════════════════════════════════════════════════════╗
# ║ Slide 14: End                                                 ║
# ╚══════════════════════════════════════════════════════════════╝
s = prs.slides.add_slide(blank)
add_bg(s, BLACK)
add_text(s, "INSIGHTX  V 4 . 0 . 0", 0.7, 0.5, 12, 0.4,
         size=11, color=INK_4, font=MONO)
add_text(s, "立即試用。", 0.7, 2.0, 12, 1.5, size=80, bold=True, color=WHITE, font=SERIF)
add_text(s, "把顧客每一句話，讀給老闆聽 — 不論是店家評論，還是 YouTube 留言。",
         0.7, 4.0, 12, 0.6, size=18, color=INK_4, font=SANS)
# CTA 區
cta = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(5.0), Inches(4.5), Inches(0.8))
cta.fill.solid(); cta.fill.fore_color.rgb = CORAL; cta.line.fill.background()
add_text(s, "github.com/GKS711/InsightX", 0.7, 5.0, 4.5, 0.8, size=15, bold=True,
         color=WHITE, font=MONO, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
add_text(s, "MIT License  ·  零瀏覽器  ·  雙平台", 5.5, 5.2, 7, 0.4,
         size=14, color=INK_3, font=SANS)
# bottom rule + footer
add_rule_line(s, 0.7, 6.8, 11.9, color=RGBColor(0x44, 0x4d, 0x49))
add_text(s, "InsightX  ·  Built with FastAPI + React 18 + Gemini  ·  Audit-friendly, evidence-backed",
         0.7, 6.95, 12, 0.3, size=10, color=INK_4, font=MONO)


# ─── Save ────────────────────────────────────────────
out = HERE / "InsightX_v4.0.0.pptx"
prs.save(str(out))
print(f"✓ Generated: {out}")
print(f"  Slides: {len(prs.slides)}")
print(f"  Size: {out.stat().st_size:,} bytes")
